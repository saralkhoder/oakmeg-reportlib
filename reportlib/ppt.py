"""Generate PowerPoint elements like tables"""
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches


def save_as_table(df: pd.DataFrame, to: str):
    """
    Write the passed DataFrame as a table in a new powerpoint file

    You should format-rename-reorder columns.

    Use functions such as
    ``
    df['ctr'].apply(lambda v: f'{v:.2f}')
    ``
    to format percentages

    Args:
        df (DataFrame): input data, can be exploded or already aggregated
        to (str): the name for the .pptx, don't write any extension here
    """
    rows, cols = df.shape
    assert (
        rows < 30 and cols < 10
    ), "Table is too large ! (more than 30 rows or 10 columns)"

    print(f"Creating table of size {rows+1}x{cols}")

    # Create presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    x, y, cx, cy = Inches(0), Inches(0), Inches(8), Inches(3)
    shape = slide.shapes.add_table(rows + 1, cols, x, y, cx, cy)
    table = shape.table

    # Write header
    for idx, col in enumerate(df.columns.str.upper()):
        table.cell(0, idx).text = col

    # Write content
    for r in range(rows):
        for c in range(cols):
            table.cell(r + 1, c).text = (
                str(df.iloc[r, c])
                if isinstance(df.iloc[r, c], str)
                else f"{df.iloc[r, c]:,}"
            )

    Path("generated").mkdir(parents=True, exist_ok=True)
    prs.save("generated/" + to + ".pptx")
